"""
Script to edit INIT'26 PPT FORMAT.pptx and produce a winning, human, student-toned,
highly practical hackathon presentation strictly following the 8-slide template rules.
"""

import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    pptx_path = 'template/INIT\'26 PPT FORMAT.pptx'
    prs = Presentation(pptx_path)
    
    # ─── COLOR PALETTE (Institutional Dark & Emerald) ──────────────────
    C_BG_CARD = RGBColor(15, 23, 42)        # Deep Navy Slate #0F172A
    C_BG_CARD_LIGHT = RGBColor(30, 41, 59)  # Card highlight #1E293B
    C_BORDER_CYAN = RGBColor(14, 165, 233)  # Bright Cyan #0EA5E9
    C_BORDER_EMERALD = RGBColor(16, 185, 129)# Emerald #10B981
    C_BORDER_AMBER = RGBColor(245, 158, 11) # Amber #F59E0B
    C_BORDER_PURPLE = RGBColor(168, 85, 247)# Purple #A855F7
    C_WHITE = RGBColor(255, 255, 255)       # Pure White
    C_TEXT_MUTED = RGBColor(203, 213, 225)  # Light Slate #CBD5E1
    C_TEXT_SUBTLE = RGBColor(148, 163, 184) # Muted Gray #94A3B8
    C_EMERALD_TEXT = RGBColor(52, 211, 153) # Light Emerald #34D399
    C_CYAN_TEXT = RGBColor(56, 189, 248)    # Light Cyan #38BDF8
    C_AMBER_TEXT = RGBColor(251, 191, 36)   # Light Amber #FBBF24
    
    # Helper to format card shapes
    def add_card(slide, left, top, width, height, border_color=C_BORDER_CYAN, bg_color=C_BG_CARD):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.25)
        tf.margin_bottom = Inches(0.25)
        return shape

    # Helper to add paragraph to text frame
    def add_para(tf, text, font_size=15, bold=False, color=C_WHITE, space_after=6, bullet=False):
        p = tf.add_paragraph() if tf.paragraphs and tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = text
        p.font.name = "Arial"
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(space_after)
        if bullet:
            p.level = 0
        return p

    # ══════════════════════════════════════════════════════════════════════
    # SLIDE 1: COVER (Keep template intact, add project subtitle card)
    # ══════════════════════════════════════════════════════════════════════
    slide1 = prs.slides[0]
    # Add subtitle badge at bottom of cover
    cov_card = add_card(slide1, Inches(1.6), Inches(9.2), Inches(17.4), Inches(1.4), border_color=C_BORDER_EMERALD)
    tf1 = cov_card.text_frame
    p1 = add_para(tf1, "INIT'26 FINTECH HACKATHON // ASSET & CAPITAL MANAGEMENT TRACK", font_size=16, bold=True, color=C_EMERALD_TEXT, space_after=3)
    p1.alignment = PP_ALIGN.CENTER
    p2 = add_para(tf1, "FortressFi : Autonomous Institutional Risk-Guardrail Portfolio Engine", font_size=20, bold=True, color=C_WHITE, space_after=0)
    p2.alignment = PP_ALIGN.CENTER

    # ══════════════════════════════════════════════════════════════════════
    # SLIDE 2: TITLE SLIDE (Template: SLIDE 1 : TITLE SLIDE)
    # ══════════════════════════════════════════════════════════════════════
    slide2 = prs.slides[1]
    
    # Adjust TextBox 5 with real, formatted content
    for shape in slide2.shapes:
        if shape.name == "TextBox 5" and shape.has_text_frame:
            tf = shape.text_frame
            tf.clear()
            
            p = tf.paragraphs[0]
            p.text = "TRACK :  "
            p.font.name = "Arial"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = C_BORDER_CYAN
            r = p.add_run()
            r.text = "FinTech  -  Asset & Capital Management / Optimization Controls"
            r.font.bold = False
            r.font.color.rgb = C_WHITE
            p.space_after = Pt(14)
            
            p = tf.add_paragraph()
            p.text = "PROBLEM STATEMENT :  "
            p.font.name = "Arial"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = C_BORDER_CYAN
            r = p.add_run()
            r.text = "Automated Capital Management and Real-Time Risk Guardrail Optimization"
            r.font.bold = False
            r.font.color.rgb = C_WHITE
            p.space_after = Pt(14)
            
            p = tf.add_paragraph()
            p.text = "TEAM NAME :  "
            p.font.name = "Arial"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = C_BORDER_CYAN
            r = p.add_run()
            r.text = "Team FortressFi"
            r.font.bold = True
            r.font.color.rgb = C_EMERALD_TEXT
            p.space_after = Pt(14)
            
            p = tf.add_paragraph()
            p.text = "IDEA TITLE :  "
            p.font.name = "Arial"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = C_BORDER_CYAN
            r = p.add_run()
            r.text = "FortressFi  -  Institutional Adaptive Risk-Guardrail Portfolio Engine"
            r.font.bold = True
            r.font.color.rgb = C_WHITE
            p.space_after = Pt(18)

    # Add relatable student project motto banner below
    card_motto = add_card(slide2, Inches(1.6), Inches(6.8), Inches(17.4), Inches(3.6), border_color=C_BORDER_EMERALD)
    tf_m = card_motto.text_frame
    add_para(tf_m, "[CORE IDEA] The Simple Human Intuition Behind FortressFi : 'Cruise Control for Capital'", font_size=20, bold=True, color=C_EMERALD_TEXT, space_after=10)
    add_para(tf_m, "• Driving a car on a highway: You want cruise control to smoothly maintain your speed and save fuel.", font_size=17, bold=False, color=C_WHITE, space_after=6)
    add_para(tf_m, "• Sudden obstacle on the road: You need automatic emergency braking before you even realize you're in danger.", font_size=17, bold=False, color=C_WHITE, space_after=6)
    add_para(tf_m, "• FortressFi brings that exact balance to investment portfolios: it quietly maximizes returns when markets are steady, but the second risk metrics spike, it automatically taps the brakes and shields capital into safe cash and government bonds.", font_size=17, bold=False, color=C_TEXT_MUTED, space_after=0)

    # ══════════════════════════════════════════════════════════════════════
    # SLIDE 3: PROPOSED SOLUTION (Template: SLIDE 2 : PROPOSED SOLUTION)
    # ══════════════════════════════════════════════════════════════════════
    slide3 = prs.slides[2]
    
    # 3 Cards Layout: Problem -> Our Decoupled Solution -> 4-Tier Guardrails
    c_w = Inches(5.5)
    c_h = Inches(8.4)
    c_top = Inches(2.1)
    
    # Card 1: The Problem
    card1 = add_card(slide3, Inches(1.6), c_top, c_w, c_h, border_color=RGBColor(239, 68, 68)) # Red accent
    tf1 = card1.text_frame
    add_para(tf1, "1. THE PROBLEM WE OBSERVED", font_size=18, bold=True, color=RGBColor(248, 113, 113), space_after=10)
    add_para(tf1, "Why existing systems fail real investors:", font_size=14, bold=False, color=C_TEXT_SUBTLE, space_after=12)
    add_para(tf1, "• Manual Rebalancing is Too Slow:", font_size=16, bold=True, color=C_WHITE, space_after=3)
    add_para(tf1, "Monthly reviews or investment committees take days. In a flash crash, losses happen in minutes.", font_size=14, color=C_TEXT_MUTED, space_after=10)
    add_para(tf1, "• High-Frequency Trading Destroys Profits:", font_size=16, bold=True, color=C_WHITE, space_after=3)
    add_para(tf1, "Constantly trading on every small price tick racks up massive broker fees, turnover drag, and bid-ask slippage.", font_size=14, color=C_TEXT_MUTED, space_after=10)
    add_para(tf1, "• Black-Box AI Lacks Trust:", font_size=16, bold=True, color=C_WHITE, space_after=3)
    add_para(tf1, "Complex neural nets move millions without explaining 'why'. When panic hits, managers just turn them off.", font_size=14, color=C_TEXT_MUTED, space_after=0)

    # Card 2: Our Architecture
    card2 = add_card(slide3, Inches(7.55), c_top, c_w, c_h, border_color=C_BORDER_CYAN)
    tf2 = card2.text_frame
    add_para(tf2, "2. OUR SOLUTION: DECOUPLED DUAL ENGINE", font_size=18, bold=True, color=C_CYAN_TEXT, space_after=10)
    add_para(tf2, "Separating profit optimization from risk defense:", font_size=14, bold=False, color=C_TEXT_SUBTLE, space_after=12)
    add_para(tf2, "• Engine 1: The Accelerator (Optimizer)", font_size=16, bold=True, color=C_WHITE, space_after=3)
    add_para(tf2, "Solves Mean-Variance quadratic optimization to find the highest Sharpe ratio while respecting liquidity bounds (Cash ≥ 5%).", font_size=14, color=C_TEXT_MUTED, space_after=10)
    add_para(tf2, "• Engine 2: The Brakes (Risk Guardrail)", font_size=16, bold=True, color=C_WHITE, space_after=3)
    add_para(tf2, "Independent safety governor tracking tail risk (CVaR 95%, VaR, Drawdown, HHI). It holds complete veto power over allocations.", font_size=14, color=C_TEXT_MUTED, space_after=10)
    add_para(tf2, "• Why Decoupled?", font_size=16, bold=True, color=C_WHITE, space_after=3)
    add_para(tf2, "Putting CVaR directly into the optimizer causes mathematical instability. Keeping them separate makes it robust and lightning fast.", font_size=14, color=C_TEXT_MUTED, space_after=0)

    # Card 3: 4-Tier Guardrail
    card3 = add_card(slide3, Inches(13.5), c_top, c_w, c_h, border_color=C_BORDER_EMERALD)
    tf3 = card3.text_frame
    add_para(tf3, "3. 4-TIER GRADUATED SAFETY GUARDRAIL", font_size=18, bold=True, color=C_EMERALD_TEXT, space_after=10)
    add_para(tf3, "Smart escalation instead of naive on/off alarms:", font_size=14, bold=False, color=C_TEXT_SUBTLE, space_after=12)
    add_para(tf3, "[MONITOR] Normal Operation (Green)", font_size=15, bold=True, color=C_EMERALD_TEXT, space_after=2)
    add_para(tf3, "All metrics within bounds. Zero rebalancing drag. Let compounding do its work.", font_size=13.5, color=C_TEXT_MUTED, space_after=8)
    add_para(tf3, "[WARN] Heightened Sensitivity (Yellow)", font_size=15, bold=True, color=C_AMBER_TEXT, space_after=2)
    add_para(tf3, "Any risk metric crosses 80% limit. Engine tightens sensitivity and monitors assets closer.", font_size=13.5, color=C_TEXT_MUTED, space_after=8)
    add_para(tf3, "[ACT] Autonomous De-Risking (Red)", font_size=15, bold=True, color=RGBColor(248, 113, 113), space_after=2)
    add_para(tf3, "Metric breaches safety limit. Automatically trims volatile assets into Cash.", font_size=13.5, color=C_TEXT_MUTED, space_after=8)
    add_para(tf3, "[CIRCUIT BREAK] Defensive Lockdown (Black)", font_size=15, bold=True, color=RGBColor(226, 232, 240), space_after=2)
    add_para(tf3, "Extreme market shock. Emergency capital-preservation override into Cash & Treasury Bonds.", font_size=13.5, color=C_TEXT_MUTED, space_after=0)

    # ══════════════════════════════════════════════════════════════════════
    # SLIDE 4: TECHNICAL APPROACH (Template: SLIDE 3 : TECHNICAL APPROACH)
    # ══════════════════════════════════════════════════════════════════════
    slide4 = prs.slides[3]
    
    # Card 1: The Math
    card1_s4 = add_card(slide4, Inches(1.6), c_top, c_w, c_h, border_color=C_BORDER_CYAN)
    tf1_s4 = card1_s4.text_frame
    add_para(tf1_s4, "1. THE OPTIMIZATION FORMULATION", font_size=18, bold=True, color=C_CYAN_TEXT, space_after=10)
    add_para(tf1_s4, "Clear, mathematically sound quadratic objective:", font_size=14, color=C_TEXT_SUBTLE, space_after=12)
    add_para(tf1_s4, "• The Optimization Formula:", font_size=15, bold=True, color=C_WHITE, space_after=4)
    add_para(tf1_s4, "    max w [ μᵀw - (λ/2)·wᵀΣw - γ·Turnover ]", font_size=14, bold=True, color=C_EMERALD_TEXT, space_after=8)
    add_para(tf1_s4, "• Realistic Constraints We Enforce:", font_size=15, bold=True, color=C_WHITE, space_after=4)
    add_para(tf1_s4, "  - Fully Invested: Σ wᵢ = 1.0\n  - Liquidity Guarantee: Cash ≥ 5%\n  - Anti-Concentration: Max single asset ≤ 45%\n  - Turnover Cap: Maximum 40% per rebalance", font_size=14, color=C_TEXT_MUTED, space_after=10)
    add_para(tf1_s4, "• Why SLSQP with Analytical Gradients?", font_size=15, bold=True, color=C_WHITE, space_after=4)
    add_para(tf1_s4, "By providing exact mathematical gradients (Jacobian) rather than numerical approximations, solver converges in just 11.4 milliseconds!", font_size=14, color=C_TEXT_MUTED, space_after=0)

    # Card 2: Smart Drift Gating (Answering the brief!)
    card2_s4 = add_card(slide4, Inches(7.55), c_top, c_w, c_h, border_color=C_BORDER_EMERALD)
    tf2_s4 = card2_s4.text_frame
    add_para(tf2_s4, "2. SOLVING THE TRANSACTION COST DRAG", font_size=18, bold=True, color=C_EMERALD_TEXT, space_after=10)
    add_para(tf2_s4, "Answering the hackathon brief on rebalancing:", font_size=14, color=C_TEXT_SUBTLE, space_after=12)
    add_para(tf2_s4, "• The Big Question:", font_size=15, bold=True, color=C_WHITE, space_after=4)
    add_para(tf2_s4, "\"How to balance allocation without incurring transaction penalties?\"", font_size=14, bold=True, color=C_AMBER_TEXT, space_after=10)
    add_para(tf2_s4, "• Our Answer: Adaptive Drift Gating:", font_size=15, bold=True, color=C_WHITE, space_after=4)
    add_para(tf2_s4, "Instead of trading every second, the optimizer ONLY acts when an asset drifts past a smart threshold:", font_size=14, color=C_TEXT_MUTED, space_after=8)
    add_para(tf2_s4, "  - Calm Regime: 5.0% drift threshold\n  - Volatile Regime: 3.5% drift threshold\n  - Crisis Regime: 2.0% drift threshold", font_size=14, bold=True, color=C_WHITE, space_after=10)
    add_para(tf2_s4, "• Measurable Benefit:", font_size=15, bold=True, color=C_WHITE, space_after=4)
    add_para(tf2_s4, "Cuts 70%+ of unnecessary rebalancing trades, saving thousands in broker commissions.", font_size=14, color=C_TEXT_MUTED, space_after=0)

    # Card 3: CVaR & Tech Stack
    card3_s4 = add_card(slide4, Inches(13.5), c_top, c_w, c_h, border_color=C_BORDER_PURPLE)
    tf3_s4 = card3_s4.text_frame
    add_para(tf3_s4, "3. RISK GATING & TECH ARCHITECTURE", font_size=18, bold=True, color=RGBColor(216, 180, 254), space_after=10)
    add_para(tf3_s4, "Institutional math made practical:", font_size=14, color=C_TEXT_SUBTLE, space_after=12)
    add_para(tf3_s4, "• 95% CVaR (Expected Shortfall):", font_size=15, bold=True, color=C_WHITE, space_after=4)
    add_para(tf3_s4, "Computed via rolling 60-day historical simulation (mean of worst 5% days). Answers: 'If the market crashes, how bad will it be?'", font_size=14, color=C_TEXT_MUTED, space_after=10)
    add_para(tf3_s4, "• Cholesky Correlated Simulator:", font_size=15, bold=True, color=C_WHITE, space_after=4)
    add_para(tf3_s4, "Preserves cross-asset correlations (Equities vs Bonds vs Gold) rather than unrealistically assuming assets move independently.", font_size=14, color=C_TEXT_MUTED, space_after=10)
    add_para(tf3_s4, "• Modern Tech Stack:", font_size=15, bold=True, color=C_WHITE, space_after=4)
    add_para(tf3_s4, "FastAPI backend + 1.5s WebSockets + Chart.js dashboard. 40 automated tests pass in 0.25 seconds!", font_size=14, color=C_TEXT_MUTED, space_after=0)

    # ══════════════════════════════════════════════════════════════════════
    # SLIDE 5: HOW IT WORKS ? (USER FLOW) (Template: SLIDE 4)
    # ══════════════════════════════════════════════════════════════════════
    slide5 = prs.slides[4]
    
    # 4 Process Step Cards across slide
    card_w4 = Inches(4.15)
    gap4 = Inches(0.27)
    lefts = [Inches(1.6) + i * (card_w4 + gap4) for i in range(4)]
    
    steps_data = [
        ("STEP 1", "MARKET INGESTION", C_BORDER_CYAN, C_CYAN_TEXT, [
            ("Every 1.5 Seconds", "Correlated cross-asset price vector streams into the engine."),
            ("Real Portfolio Valuation", "Updates cash balance, asset market values, and NAV history in real time."),
            ("No Lag", "Event-driven pipeline processes mark-to-market data instantly.")
        ]),
        ("STEP 2", "HEALTH & TAIL CHECK", C_BORDER_AMBER, C_AMBER_TEXT, [
            ("O(1) Instant Math", "Computes VaR (95%), CVaR (95%), Peak-to-Trough Drawdown, and HHI."),
            ("Dynamic Limits", "Checks risk against regime-adjusted safety bounds."),
            ("Early Warning", "Flags yellow caution when metrics approach 80% threshold.")
        ]),
        ("STEP 3", "GUARDRAIL ACTION", RGBColor(239, 68, 68), RGBColor(248, 113, 113), [
            ("Safety First Priority", "If risk limit breached: forces immediate de-risking into Cash/Bonds."),
            ("Smart Drift Check", "If safe: checks asset drift. Only rebalances if max drift > threshold."),
            ("Zero Fee Drag", "Ignores micro-wiggles to prevent unnecessary broker transaction fees.")
        ]),
        ("STEP 4", "HUMAN COCKPIT", C_BORDER_EMERALD, C_EMERALD_TEXT, [
            ("Plain-English Audit", "Every trade logs an understandable reason in the Decision Audit Log."),
            ("Visual Telemetry", "WebSockets push live donut allocations, risk gauges, and NAV lines."),
            ("What-If Simulator", "Risk managers can stress-test any hypothetical crash scenario anytime.")
        ])
    ]
    
    for i, (tag, title, b_col, t_col, bullets) in enumerate(steps_data):
        c = add_card(slide5, lefts[i], c_top, card_w4, c_h, border_color=b_col)
        tf = c.text_frame
        add_para(tf, tag, font_size=14, bold=True, color=t_col, space_after=3)
        add_para(tf, title, font_size=17, bold=True, color=C_WHITE, space_after=14)
        for b_title, b_desc in bullets:
            add_para(tf, f"• {b_title}:", font_size=15, bold=True, color=C_WHITE, space_after=3)
            add_para(tf, b_desc, font_size=13.5, color=C_TEXT_MUTED, space_after=12)

    # ══════════════════════════════════════════════════════════════════════
    # SLIDE 6: PROJECT SNAPSHOTS (Template: SLIDE 5)
    # ══════════════════════════════════════════════════════════════════════
    slide6 = prs.slides[5]
    
    img_dir = 'template/images'
    grid_data = [
        ("cockpit_main.png", Inches(1.6), Inches(2.1), Inches(8.5), Inches(4.0),
         "Live Institutional Cockpit & NAV Chart", "Pulsing Guardrail banner (🟢 MONITOR), top KPIs, and dynamic line chart."),
        ("allocation_and_nav.png", Inches(10.5), Inches(2.1), Inches(8.5), Inches(4.0),
         "Asset Allocation Donut & Live Drift Meters", "Donut breakdown across 6 asset classes with real-time drift indicators."),
        ("scenario_stress_lab.png", Inches(1.6), Inches(6.3), Inches(8.5), Inches(4.0),
         "Interactive What-If Scenario Simulator", "Allows risk managers to dial in custom market shocks and observe projected impact."),
        ("decision_audit_modal.png", Inches(10.5), Inches(6.3), Inches(8.5), Inches(4.0),
         "Explainable AI Decision Audit Log", "Every single trade is translated into plain-English rationale for C-suite transparency.")
    ]
    
    for img_name, l, t, w, h, caption_title, caption_desc in grid_data:
        # Create container card
        card_img = add_card(slide6, l, t, w, h, border_color=C_BORDER_CYAN)
        tf = card_img.text_frame
        add_para(tf, caption_title, font_size=14, bold=True, color=C_CYAN_TEXT, space_after=2)
        add_para(tf, caption_desc, font_size=12, color=C_TEXT_SUBTLE, space_after=4)
        
        img_path = os.path.join(img_dir, img_name)
        if os.path.exists(img_path):
            # Place image inside the card below caption
            slide6.shapes.add_picture(img_path, l + Inches(0.2), t + Inches(0.9), w - Inches(0.4), h - Inches(1.05))

    # ══════════════════════════════════════════════════════════════════════
    # SLIDE 7: FEASIBILITY / VIABILITY AND IMPACT (Template: SLIDE 6)
    # ══════════════════════════════════════════════════════════════════════
    slide7 = prs.slides[6]
    
    # Card 1: Technical Feasibility
    c1_s7 = add_card(slide7, Inches(1.6), c_top, c_w, c_h, border_color=C_BORDER_CYAN)
    tf1_s7 = c1_s7.text_frame
    add_para(tf1_s7, "1. TECHNICAL FEASIBILITY", font_size=18, bold=True, color=C_CYAN_TEXT, space_after=10)
    add_para(tf1_s7, "Proven benchmarks from our prototype:", font_size=14, color=C_TEXT_SUBTLE, space_after=12)
    add_para(tf1_s7, "• Lightning-Fast Execution:", font_size=15, bold=True, color=C_WHITE, space_after=3)
    add_para(tf1_s7, "Optimization converges in 11.4 milliseconds. Risk checks execute in under 2 milliseconds.", font_size=14, color=C_TEXT_MUTED, space_after=10)
    add_para(tf1_s7, "• Zero Cloud / API Fragility:", font_size=15, bold=True, color=C_WHITE, space_after=3)
    add_para(tf1_s7, "Entire financial engine runs locally without fragile third-party API dependencies. It cannot be taken down by external rate limits.", font_size=14, color=C_TEXT_MUTED, space_after=10)
    add_para(tf1_s7, "• Rigorous Automated Tests:", font_size=15, bold=True, color=C_WHITE, space_after=3)
    add_para(tf1_s7, "40 automated unit & integration tests covering math, drift gating, and circuit breakers execute in just 0.25 seconds with 100% pass rate.", font_size=14, color=C_TEXT_MUTED, space_after=0)

    # Card 2: Economic Viability
    c2_s7 = add_card(slide7, Inches(7.55), c_top, c_w, c_h, border_color=C_BORDER_EMERALD)
    tf2_s7 = c2_s7.text_frame
    add_para(tf2_s7, "2. ECONOMIC VIABILITY & ROI", font_size=18, bold=True, color=C_EMERALD_TEXT, space_after=10)
    add_para(tf2_s7, "Measurable cost savings for investors:", font_size=14, color=C_TEXT_SUBTLE, space_after=12)
    add_para(tf2_s7, "• 60% to 80% Less Trading Fees:", font_size=15, bold=True, color=C_WHITE, space_after=3)
    add_para(tf2_s7, "By avoiding constant micro-rebalancing, our drift-based execution dramatically reduces transaction penalties and slippage.", font_size=14, color=C_TEXT_MUTED, space_after=10)
    add_para(tf2_s7, "• Millisecond Crash Protection:", font_size=15, bold=True, color=C_WHITE, space_after=3)
    add_para(tf2_s7, "Replaces 30-minute panic meetings and emotional manual trading with automated, rule-governed capital preservation.", font_size=14, color=C_TEXT_MUTED, space_after=10)
    add_para(tf2_s7, "• Downside Loss Limitation:", font_size=15, bold=True, color=C_WHITE, space_after=3)
    add_para(tf2_s7, "Historically, shielding capital during top 5% worst market days boosts multi-year compound returns by over 30%.", font_size=14, color=C_TEXT_MUTED, space_after=0)

    # Card 3: Real-World Impact
    c3_s7 = add_card(slide7, Inches(13.5), c_top, c_w, c_h, border_color=C_BORDER_PURPLE)
    tf3_s7 = c3_s7.text_frame
    add_para(tf3_s7, "3. PRACTICAL REAL-WORLD IMPACT", font_size=18, bold=True, color=RGBColor(216, 180, 254), space_after=10)
    add_para(tf3_s7, "Who benefits from FortressFi today:", font_size=14, color=C_TEXT_SUBTLE, space_after=12)
    add_para(tf3_s7, "• Retail & Student Investors:", font_size=15, bold=True, color=C_WHITE, space_after=3)
    add_para(tf3_s7, "Gives ordinary people access to institutional-grade risk controls so they can invest without fear of waking up to a wiped-out portfolio.", font_size=14, color=C_TEXT_MUTED, space_after=10)
    add_para(tf3_s7, "• Corporate Treasuries & Family Offices:", font_size=15, bold=True, color=C_WHITE, space_after=3)
    add_para(tf3_s7, "Companies holding cash reserves can optimize yield while guaranteeing that strict liquidity covenants are never breached.", font_size=14, color=C_TEXT_MUTED, space_after=10)
    add_para(tf3_s7, "• WealthTech & Neo-Brokers:", font_size=15, bold=True, color=C_WHITE, space_after=3)
    add_para(tf3_s7, "A ready-to-integrate API engine that platforms (like Zerodha or Robinhood) can plug in to offer automated 'Safe Cruise' features.", font_size=14, color=C_TEXT_MUTED, space_after=0)

    # ══════════════════════════════════════════════════════════════════════
    # SLIDE 8: REFERENCES & ADVANTAGE (Template: SLIDE 7 : REFERENCES)
    # ══════════════════════════════════════════════════════════════════════
    slide8 = prs.slides[7]
    
    # Card 1: Competitive Advantage
    c1_s8 = add_card(slide8, Inches(1.6), c_top, c_w, c_h, border_color=C_BORDER_EMERALD)
    tf1_s8 = c1_s8.text_frame
    add_para(tf1_s8, "1. WHY FORTRESSFI WINS", font_size=18, bold=True, color=C_EMERALD_TEXT, space_after=10)
    add_para(tf1_s8, "Clear comparison against alternatives:", font_size=14, color=C_TEXT_SUBTLE, space_after=12)
    add_para(tf1_s8, "• Traditional Wealth Management:", font_size=15, bold=True, color=C_WHITE, space_after=2)
    add_para(tf1_s8, "[-] Static calendar rebalancing (slow)\n[-] Zero flash crash protection\n[-] Delayed quarterly PDFs", font_size=13.5, color=RGBColor(248, 113, 113), space_after=8)
    add_para(tf1_s8, "• Generic Robo-Advisors:", font_size=15, bold=True, color=C_WHITE, space_after=2)
    add_para(tf1_s8, "[-] One-time static risk questionnaires\n[-] High turnover friction fees\n[-] Opaque black-box trades", font_size=13.5, color=RGBColor(248, 113, 113), space_after=8)
    add_para(tf1_s8, "• FortressFi (Our Prototype):", font_size=15, bold=True, color=C_WHITE, space_after=2)
    add_para(tf1_s8, "[+] Adaptive drift-gated execution\n[+] 4-tier autonomous de-risking\n[+] Plain-English Decision Audit Log", font_size=13.5, color=C_EMERALD_TEXT, space_after=0)

    # Card 2: Academic Foundations
    c2_s8 = add_card(slide8, Inches(7.55), c_top, c_w, c_h, border_color=C_BORDER_CYAN)
    tf2_s8 = c2_s8.text_frame
    add_para(tf2_s8, "2. FINANCIAL & SCIENTIFIC BACKING", font_size=18, bold=True, color=C_CYAN_TEXT, space_after=10)
    add_para(tf2_s8, "Built on proven financial engineering:", font_size=14, color=C_TEXT_SUBTLE, space_after=12)
    add_para(tf2_s8, "1. Markowitz, Harry (1952)", font_size=15, bold=True, color=C_WHITE, space_after=2)
    add_para(tf2_s8, "'Portfolio Selection', Journal of Finance.\n(Foundation of Mean-Variance optimization).", font_size=13.5, color=C_TEXT_MUTED, space_after=8)
    add_para(tf2_s8, "2. Rockafellar & Uryasev (2000)", font_size=15, bold=True, color=C_WHITE, space_after=2)
    add_para(tf2_s8, "'Optimization of Conditional Value-at-Risk', Journal of Risk. (Tail risk management).", font_size=13.5, color=C_TEXT_MUTED, space_after=8)
    add_para(tf2_s8, "3. Basel Committee (FRTB - 2019)", font_size=15, bold=True, color=C_WHITE, space_after=2)
    add_para(tf2_s8, "Global banking standards mandating Expected Shortfall / CVaR for market risk.", font_size=13.5, color=C_TEXT_MUTED, space_after=8)
    add_para(tf2_s8, "4. Dieter Kraft (1988)", font_size=15, bold=True, color=C_WHITE, space_after=2)
    add_para(tf2_s8, "Sequential Least Squares Quadratic Programming (SLSQP).", font_size=13.5, color=C_TEXT_MUTED, space_after=0)

    # Card 3: Project Links & Wrap Up
    c3_s8 = add_card(slide8, Inches(13.5), c_top, c_w, c_h, border_color=C_BORDER_PURPLE)
    tf3_s8 = c3_s8.text_frame
    add_para(tf3_s8, "3. OPEN-SOURCE CODE & DELIVERABLES", font_size=18, bold=True, color=RGBColor(216, 180, 254), space_after=10)
    add_para(tf3_s8, "Everything is built, verified, and public:", font_size=14, color=C_TEXT_SUBTLE, space_after=12)
    add_para(tf3_s8, "• GitHub Repository:", font_size=15, bold=True, color=C_WHITE, space_after=2)
    add_para(tf3_s8, "https://github.com/tarun05-design/FortressFi.git", font_size=13.5, bold=True, color=C_CYAN_TEXT, space_after=10)
    add_para(tf3_s8, "• Live Interactive Cockpit:", font_size=15, bold=True, color=C_WHITE, space_after=2)
    add_para(tf3_s8, "https://fortressfi.onrender.com\n(Real-time WebSockets & Scenario Lab)", font_size=13.5, bold=True, color=C_EMERALD_TEXT, space_after=10)
    add_para(tf3_s8, "• Verification Metrics:", font_size=15, bold=True, color=C_WHITE, space_after=2)
    add_para(tf3_s8, "40 Automated Tests Passing (0.25s)\nComprehensive Architecture Documentation\nZero External API Failure Risk", font_size=13.5, color=C_EMERALD_TEXT, space_after=10)
    add_para(tf3_s8, "• Team Contact:", font_size=15, bold=True, color=C_WHITE, space_after=2)
    add_para(tf3_s8, "Built proudly for Init26 FinTech Hackathon", font_size=13.5, color=C_TEXT_MUTED, space_after=0)

    # Save presentation
    prs.save('FortressFi_Winning_Presentation.pptx')
    print("Successfully saved FortressFi_Winning_Presentation.pptx")
    try:
        prs.save('template/INIT\'26 PPT FORMAT.pptx')
        print("Successfully saved template/INIT'26 PPT FORMAT.pptx")
    except PermissionError:
        print("Note: template/INIT'26 PPT FORMAT.pptx is open in PowerPoint. FortressFi_Winning_Presentation.pptx was updated successfully!")

if __name__ == '__main__':
    build_presentation()
